"""PowerPoint document converter implementation."""

import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Union

import structlog
from pptx import Presentation

from morag_core.interfaces.converter import (
    ConversionOptions,
    ConversionError,
)
from morag_core.models.document import Document, DocumentType

from .base import DocumentConverter

logger = structlog.get_logger(__name__)


class PresentationConverter(DocumentConverter):
    """PowerPoint document converter implementation."""

    def __init__(self):
        """Initialize PowerPoint converter."""
        super().__init__()
        self.supported_formats = {"powerpoint", "pptx", "ppt"}

    async def _extract_text(self, file_path: Path, document: Document, options: ConversionOptions) -> Document:
        """Extract text from PowerPoint document.

        Args:
            file_path: Path to PowerPoint file
            document: Document to update
            options: Conversion options

        Returns:
            Updated document

        Raises:
            ConversionError: If text extraction fails
        """
        try:
            # Load presentation
            presentation = Presentation(file_path)
            
            # Extract metadata
            if options.extract_metadata:
                document.metadata.title = os.path.basename(file_path)
                document.metadata.file_type = "powerpoint"
                document.metadata.page_count = len(presentation.slides)
                
                # Try to get properties if available
                try:
                    core_props = presentation.core_properties
                    if core_props.title:
                        document.metadata.title = core_props.title
                    if core_props.author:
                        document.metadata.author = core_props.author
                    if core_props.created:
                        document.metadata.creation_date = core_props.created.isoformat()
                    if core_props.modified:
                        document.metadata.modification_date = core_props.modified.isoformat()
                except Exception as e:
                    logger.warning("Failed to extract PowerPoint properties", error=str(e))
            
            # Process based on chunking strategy
            if options.chunking_strategy == "slide":
                # Process each slide as a separate chunk
                for i, slide in enumerate(presentation.slides):
                    slide_text = self._extract_slide_text(slide)
                    slide_number = i + 1
                    
                    # Add slide as chunk
                    document.add_chunk(
                        content=slide_text,
                        section=f"Slide {slide_number}",
                    )
            else:
                # Process all slides as a single document
                all_text = []
                for i, slide in enumerate(presentation.slides):
                    slide_text = self._extract_slide_text(slide)
                    slide_number = i + 1
                    all_text.append(f"Slide {slide_number}:\n\n{slide_text}")
                
                # Join all text
                document.raw_text = "\n\n".join(all_text)
                
                # Apply default chunking strategy
                await self._chunk_document(document, options)
            
            # Estimate word count
            document.metadata.word_count = len(document.raw_text.split())
            
            return document
            
        except Exception as e:
            logger.error(
                "PowerPoint document extraction failed",
                error=str(e),
                error_type=e.__class__.__name__,
                file_path=str(file_path),
            )
            raise ConversionError(f"Failed to extract text from PowerPoint document: {str(e)}")
    
    def _extract_slide_text(self, slide) -> str:
        """Extract text from PowerPoint slide.

        Args:
            slide: PowerPoint slide

        Returns:
            Extracted text
        """
        slide_text = []
        
        # Extract title if available
        if slide.shapes.title and slide.shapes.title.text:
            slide_text.append(f"Title: {slide.shapes.title.text}")
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                # Skip empty text
                if shape.text.strip():
                    slide_text.append(shape.text)
            
            # Extract text from tables
            if shape.has_table:
                table_text = []
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                if table_text:
                    slide_text.append("Table:\n" + "\n".join(table_text))
        
        return "\n\n".join(slide_text)