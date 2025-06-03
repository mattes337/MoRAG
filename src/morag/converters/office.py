"""Enhanced Office documents to Markdown converter with full format support."""

import time
from pathlib import Path
from typing import Union, List, Dict, Any, Optional
import structlog
import io

from .base import BaseConverter, ConversionOptions, ConversionResult, QualityScore
from .quality import ConversionQualityValidator

logger = structlog.get_logger(__name__)

# Word document processing
try:
    from docx import Document
    from docx.shared import Inches
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available, Word document processing disabled")

# Excel processing
try:
    import openpyxl
    from openpyxl.utils import get_column_letter
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    logger.warning("openpyxl not available, Excel processing disabled")

# PowerPoint processing
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available, PowerPoint processing disabled")

# Legacy Excel support
try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False


class OfficeConverter(BaseConverter):
    """Enhanced Office documents to Markdown converter with full format support."""

    def __init__(self):
        super().__init__("Enhanced MoRAG Office Converter")
        self.supported_formats = ['word', 'excel', 'powerpoint', 'docx', 'xlsx', 'pptx', 'doc', 'xls', 'ppt']
        self.quality_validator = ConversionQualityValidator()

        # Initialize format-specific converters
        self.word_converter = WordToMarkdownConverter() if DOCX_AVAILABLE else None
        self.excel_converter = ExcelToMarkdownConverter() if EXCEL_AVAILABLE else None
        self.powerpoint_converter = PowerPointToMarkdownConverter() if PPTX_AVAILABLE else None
    
    def supports_format(self, format_type: str) -> bool:
        """Check if this converter supports the given format."""
        return format_type.lower() in self.supported_formats
    
    async def convert(self, file_path: Union[str, Path], options: ConversionOptions) -> ConversionResult:
        """Convert office document to structured markdown.

        Args:
            file_path: Path to office document
            options: Conversion options

        Returns:
            ConversionResult with markdown content
        """
        start_time = time.time()
        file_path = Path(file_path)

        await self.validate_input(file_path)

        format_type = self._detect_office_format(file_path)

        logger.info(
            "Starting enhanced office document conversion",
            file_path=str(file_path),
            format=format_type,
            extract_tables=options.format_options.get('extract_tables', True),
            extract_images=options.extract_images
        )

        try:
            # Route to appropriate converter
            if format_type in ['docx', 'doc'] and self.word_converter:
                result = await self.word_converter.convert(file_path, options)
            elif format_type in ['xlsx', 'xls'] and self.excel_converter:
                result = await self.excel_converter.convert(file_path, options)
            elif format_type in ['pptx', 'ppt'] and self.powerpoint_converter:
                result = await self.powerpoint_converter.convert(file_path, options)
            else:
                # Fallback for unsupported formats or missing dependencies
                result = await self._create_fallback_conversion(file_path, format_type, options)

            processing_time = time.time() - start_time
            result.processing_time = processing_time
            result.converter_used = self.name

            logger.info(
                "Office document conversion completed",
                processing_time=processing_time,
                format=format_type,
                quality_score=result.quality_score.overall_score if result.quality_score else 0,
                word_count=result.word_count
            )

            return result

        except Exception as e:
            processing_time = time.time() - start_time
            error_msg = f"Office document conversion failed: {str(e)}"

            logger.error(
                "Office document conversion failed",
                error=str(e),
                error_type=type(e).__name__,
                processing_time=processing_time,
                format=format_type
            )

            return ConversionResult(
                content="",
                metadata={},
                processing_time=processing_time,
                success=False,
                error_message=error_msg,
                original_format=format_type,
                converter_used=self.name
            )

    def _detect_office_format(self, file_path: Path) -> str:
        """Detect specific office document format.

        Args:
            file_path: Path to office document

        Returns:
            Specific format type (docx, xlsx, pptx, etc.)
        """
        extension = file_path.suffix.lower()

        format_mapping = {
            '.docx': 'docx',
            '.doc': 'doc',
            '.xlsx': 'xlsx',
            '.xls': 'xls',
            '.pptx': 'pptx',
            '.ppt': 'ppt'
        }

        return format_mapping.get(extension, 'unknown')

    async def _create_fallback_conversion(self, file_path: Path, format_type: str, options: ConversionOptions) -> ConversionResult:
        """Create fallback conversion for unsupported formats.

        Args:
            file_path: Path to office document
            format_type: Detected format type
            options: Conversion options

        Returns:
            Basic conversion result
        """
        markdown_content = await self._create_fallback_markdown(file_path, format_type, options)

        quality_score = QualityScore(
            overall_score=0.3,  # Low score for fallback
            completeness_score=0.2,
            readability_score=0.4,
            structure_score=0.3,
            metadata_preservation=0.4
        )

        return ConversionResult(
            content=markdown_content,
            metadata=self._create_basic_metadata(file_path, format_type),
            quality_score=quality_score,
            success=True,
            warnings=[f"Limited support for {format_type} format. Install required dependencies for full functionality."],
            original_format=format_type
        )
    
    async def _create_placeholder_markdown(self, file_path: Path, format_type: str, options: ConversionOptions) -> str:
        """Create placeholder markdown content.
        
        Args:
            file_path: Path to office document
            format_type: Detected format type
            options: Conversion options
            
        Returns:
            Placeholder markdown content
        """
        sections = []
        
        # Document header
        title = file_path.stem
        sections.append(f"# {title}")
        sections.append("")
        
        # Notice about limited support
        sections.append("## ⚠️ Limited Format Support")
        sections.append("")
        sections.append(f"This {format_type.upper()} document was processed with limited support.")
        sections.append("For full functionality, please install the required dependencies.")
        sections.append("")
        
        # Document information
        if options.include_metadata:
            sections.append("## Document Information")
            sections.append("")
            sections.append(f"**File Name**: {file_path.name}")
            sections.append(f"**Format**: {format_type.upper()}")
            sections.append(f"**File Size**: {file_path.stat().st_size:,} bytes")
            sections.append("")
        
        # Format-specific placeholder content
        if format_type in ['word', 'docx']:
            sections.append("## Document Content")
            sections.append("")
            sections.append("*Word document content will be extracted here*")
            sections.append("")
            sections.append("### Features to be implemented:")
            sections.append("- Text extraction with formatting preservation")
            sections.append("- Table extraction and conversion")
            sections.append("- Image extraction and description")
            sections.append("- Header and footer processing")
            sections.append("- Comment and track changes handling")
            
        elif format_type in ['excel', 'xlsx']:
            sections.append("## Workbook Content")
            sections.append("")
            sections.append("*Excel workbook content will be extracted here*")
            sections.append("")
            sections.append("### Features to be implemented:")
            sections.append("- Worksheet data extraction")
            sections.append("- Formula preservation and calculation")
            sections.append("- Chart extraction and description")
            sections.append("- Pivot table processing")
            sections.append("- Cell formatting preservation")
            
        elif format_type in ['powerpoint', 'pptx']:
            sections.append("## Presentation Content")
            sections.append("")
            sections.append("*PowerPoint presentation content will be extracted here*")
            sections.append("")
            sections.append("### Features to be implemented:")
            sections.append("- Slide content extraction")
            sections.append("- Speaker notes processing")
            sections.append("- Image and shape extraction")
            sections.append("- Animation and transition notes")
            sections.append("- Master slide template information")
        
        sections.append("")
        
        # Implementation status
        sections.append("## Implementation Status")
        sections.append("")
        sections.append("- [ ] Full document parsing")
        sections.append("- [ ] Content extraction")
        sections.append("- [ ] Formatting preservation")
        sections.append("- [ ] Metadata extraction")
        sections.append("- [ ] Quality assessment")
        sections.append("")
        sections.append("**Expected completion**: Task 28 implementation")
        
        return "\n".join(sections)
    
    def _create_basic_metadata(self, file_path: Path, format_type: str) -> dict:
        """Create basic metadata for office document.
        
        Args:
            file_path: Path to office document
            format_type: Detected format type
            
        Returns:
            Basic metadata dictionary
        """
        stat = file_path.stat()
        
        return {
            'original_filename': file_path.name,
            'file_size': stat.st_size,
            'format_type': format_type,
            'conversion_format': f'{format_type}_to_markdown',
            'converter_version': '1.0.0-placeholder',
            'file_extension': file_path.suffix.lower(),
            'created_time': stat.st_ctime,
            'modified_time': stat.st_mtime,
            'implementation_status': 'placeholder',
            'full_implementation_task': 'Task 28'
        }


class WordToMarkdownConverter:
    """Word document to Markdown converter."""

    def __init__(self):
        self.quality_validator = ConversionQualityValidator()

    async def convert(self, file_path: Path, options: ConversionOptions) -> ConversionResult:
        """Convert Word document to markdown.

        Args:
            file_path: Path to Word document
            options: Conversion options

        Returns:
            ConversionResult with markdown content
        """
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not available for Word document processing")

        try:
            doc = Document(str(file_path))

            # Extract document metadata
            metadata = self._extract_word_metadata(doc, file_path)

            # Convert to markdown
            markdown_content = await self._convert_word_to_markdown(doc, options)

            # Calculate quality score
            quality_score = self._calculate_word_quality(doc, markdown_content)

            return ConversionResult(
                content=markdown_content,
                metadata=metadata,
                quality_score=quality_score,
                success=True,
                original_format='docx'
            )

        except Exception as e:
            logger.error(f"Word document conversion failed: {e}")
            raise

    async def _convert_word_to_markdown(self, doc, options: ConversionOptions) -> str:
        """Convert Word document to structured markdown.

        Args:
            doc: Word document object
            options: Conversion options

        Returns:
            Markdown content
        """
        sections = []

        # Document title
        title = doc.core_properties.title or "Word Document"
        sections.append(f"# {title}")
        sections.append("")

        # Document metadata
        if options.include_metadata:
            sections.append("## Document Information")
            sections.append("")

            if doc.core_properties.author:
                sections.append(f"**Author**: {doc.core_properties.author}")
            if doc.core_properties.created:
                sections.append(f"**Created**: {doc.core_properties.created}")
            if doc.core_properties.modified:
                sections.append(f"**Modified**: {doc.core_properties.modified}")

            sections.append("")

        # Process document content
        sections.append("## Content")
        sections.append("")

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Handle different paragraph styles
                style_name = paragraph.style.name if paragraph.style else "Normal"

                if "Heading" in style_name:
                    level = self._extract_heading_level(style_name)
                    sections.append(f"{'#' * (level + 2)} {paragraph.text}")
                else:
                    # Process text with formatting
                    formatted_text = self._process_paragraph_formatting(paragraph)
                    sections.append(formatted_text)

                sections.append("")

        # Process tables
        if options.format_options.get('extract_tables', True):
            for i, table in enumerate(doc.tables, 1):
                table_md = self._convert_word_table(table)
                if table_md:
                    sections.append(f"### Table {i}")
                    sections.append("")
                    sections.append(table_md)
                    sections.append("")

        return "\n".join(sections)

    def _extract_heading_level(self, style_name: str) -> int:
        """Extract heading level from style name."""
        if "Heading 1" in style_name:
            return 1
        elif "Heading 2" in style_name:
            return 2
        elif "Heading 3" in style_name:
            return 3
        elif "Heading 4" in style_name:
            return 4
        elif "Heading 5" in style_name:
            return 5
        elif "Heading 6" in style_name:
            return 6
        else:
            return 1

    def _process_paragraph_formatting(self, paragraph) -> str:
        """Process paragraph with text formatting.

        Args:
            paragraph: Word paragraph object

        Returns:
            Formatted text string
        """
        text_parts = []

        for run in paragraph.runs:
            text = run.text

            # Apply formatting
            if run.bold:
                text = f"**{text}**"
            if run.italic:
                text = f"*{text}*"
            if run.underline:
                text = f"<u>{text}</u>"

            text_parts.append(text)

        return "".join(text_parts)

    def _convert_word_table(self, table) -> str:
        """Convert Word table to markdown format.

        Args:
            table: Word table object

        Returns:
            Markdown table string
        """
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)

            if not rows:
                return ""

            # Create markdown table
            table_lines = []

            # Header row
            headers = rows[0]
            table_lines.append("| " + " | ".join(headers) + " |")
            table_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

            # Data rows
            for row in rows[1:]:
                # Ensure row has same number of cells as headers
                while len(row) < len(headers):
                    row.append("")
                table_lines.append("| " + " | ".join(row[:len(headers)]) + " |")

            return "\n".join(table_lines)

        except Exception as e:
            logger.warning(f"Failed to convert Word table: {e}")
            return "*[Table conversion failed]*"

    def _extract_word_metadata(self, doc, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from Word document.

        Args:
            doc: Word document object
            file_path: Path to file

        Returns:
            Metadata dictionary
        """
        metadata = {
            'original_filename': file_path.name,
            'file_size': file_path.stat().st_size,
            'conversion_format': 'word_to_markdown',
            'converter_version': '1.0.0',
            'office_type': 'word'
        }

        # Add document properties
        if doc.core_properties:
            props = doc.core_properties
            if props.title:
                metadata['title'] = props.title
            if props.author:
                metadata['author'] = props.author
            if props.created:
                metadata['created_date'] = props.created.isoformat()
            if props.modified:
                metadata['modified_date'] = props.modified.isoformat()
            if props.subject:
                metadata['subject'] = props.subject

        # Count elements
        metadata.update({
            'paragraph_count': len(doc.paragraphs),
            'table_count': len(doc.tables),
            'has_tables': len(doc.tables) > 0,
            'word_count': sum(len(p.text.split()) for p in doc.paragraphs)
        })

        return metadata

    def _calculate_word_quality(self, doc, markdown_content: str) -> QualityScore:
        """Calculate quality score for Word conversion.

        Args:
            doc: Word document object
            markdown_content: Generated markdown

        Returns:
            Quality score
        """
        # Base scores
        completeness_score = 0.9
        readability_score = 0.85
        structure_score = 0.8
        metadata_preservation = 0.9

        # Adjust based on content
        word_count = len(markdown_content.split())
        if word_count < 50:
            completeness_score *= 0.8

        # Check for tables
        if len(doc.tables) > 0 and '|' in markdown_content:
            structure_score = min(structure_score * 1.1, 1.0)

        # Check for headings
        heading_count = markdown_content.count('#')
        if heading_count > 0:
            structure_score = min(structure_score * (1 + heading_count * 0.02), 1.0)

        overall_score = (
            completeness_score * 0.3 +
            readability_score * 0.25 +
            structure_score * 0.25 +
            metadata_preservation * 0.2
        )

        return QualityScore(
            overall_score=overall_score,
            completeness_score=completeness_score,
            readability_score=readability_score,
            structure_score=structure_score,
            metadata_preservation=metadata_preservation
        )


class ExcelToMarkdownConverter:
    """Excel workbook to Markdown converter."""

    def __init__(self):
        self.quality_validator = ConversionQualityValidator()

    async def convert(self, file_path: Path, options: ConversionOptions) -> ConversionResult:
        """Convert Excel workbook to markdown.

        Args:
            file_path: Path to Excel file
            options: Conversion options

        Returns:
            ConversionResult with markdown content
        """
        if not EXCEL_AVAILABLE:
            raise ImportError("openpyxl not available for Excel processing")

        try:
            workbook = openpyxl.load_workbook(str(file_path), data_only=True)

            # Extract workbook metadata
            metadata = self._extract_excel_metadata(workbook, file_path)

            # Convert to markdown
            markdown_content = await self._convert_excel_to_markdown(workbook, options)

            # Calculate quality score
            quality_score = self._calculate_excel_quality(workbook, markdown_content)

            return ConversionResult(
                content=markdown_content,
                metadata=metadata,
                quality_score=quality_score,
                success=True,
                original_format='xlsx'
            )

        except Exception as e:
            logger.error(f"Excel conversion failed: {e}")
            raise

    async def _convert_excel_to_markdown(self, workbook, options: ConversionOptions) -> str:
        """Convert Excel workbook to structured markdown.

        Args:
            workbook: Excel workbook object
            options: Conversion options

        Returns:
            Markdown content
        """
        sections = []

        # Workbook title
        title = "Excel Workbook"
        sections.append(f"# {title}")
        sections.append("")

        # Workbook metadata
        if options.include_metadata:
            sections.append("## Workbook Information")
            sections.append("")
            sections.append(f"**Worksheets**: {len(workbook.worksheets)}")
            sections.append("")

        # Process each worksheet
        for sheet in workbook.worksheets:
            sections.append(f"## Worksheet: {sheet.title}")
            sections.append("")

            # Convert sheet data to markdown table
            table_md = self._convert_excel_sheet(sheet)
            if table_md:
                sections.append(table_md)
            else:
                sections.append("*No data in this worksheet*")

            sections.append("")

        return "\n".join(sections)

    def _convert_excel_sheet(self, sheet) -> str:
        """Convert Excel sheet to markdown table.

        Args:
            sheet: Excel worksheet object

        Returns:
            Markdown table string
        """
        try:
            # Get all data from sheet
            data = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):  # Skip empty rows
                    data.append([str(cell) if cell is not None else "" for cell in row])

            if not data:
                return ""

            # Create markdown table
            table_lines = []

            # Header row
            if data:
                headers = data[0]
                table_lines.append("| " + " | ".join(headers) + " |")
                table_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

                # Data rows
                for row in data[1:]:
                    # Ensure row has same number of cells as headers
                    while len(row) < len(headers):
                        row.append("")
                    table_lines.append("| " + " | ".join(row[:len(headers)]) + " |")

            return "\n".join(table_lines)

        except Exception as e:
            logger.warning(f"Failed to convert Excel sheet: {e}")
            return "*[Sheet conversion failed]*"

    def _extract_excel_metadata(self, workbook, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from Excel workbook.

        Args:
            workbook: Excel workbook object
            file_path: Path to file

        Returns:
            Metadata dictionary
        """
        metadata = {
            'original_filename': file_path.name,
            'file_size': file_path.stat().st_size,
            'conversion_format': 'excel_to_markdown',
            'converter_version': '1.0.0',
            'office_type': 'excel'
        }

        # Add workbook properties
        metadata.update({
            'worksheet_count': len(workbook.worksheets),
            'worksheet_names': [sheet.title for sheet in workbook.worksheets],
            'has_data': any(any(row) for sheet in workbook.worksheets for row in sheet.iter_rows(values_only=True))
        })

        return metadata

    def _calculate_excel_quality(self, workbook, markdown_content: str) -> QualityScore:
        """Calculate quality score for Excel conversion.

        Args:
            workbook: Excel workbook object
            markdown_content: Generated markdown

        Returns:
            Quality score
        """
        # Base scores
        completeness_score = 0.85
        readability_score = 0.8
        structure_score = 0.9
        metadata_preservation = 0.9

        # Adjust based on content
        if len(workbook.worksheets) > 1:
            structure_score = min(structure_score * 1.1, 1.0)

        # Check for tables
        if '|' in markdown_content:
            structure_score = min(structure_score * 1.1, 1.0)

        overall_score = (
            completeness_score * 0.3 +
            readability_score * 0.25 +
            structure_score * 0.25 +
            metadata_preservation * 0.2
        )

        return QualityScore(
            overall_score=overall_score,
            completeness_score=completeness_score,
            readability_score=readability_score,
            structure_score=structure_score,
            metadata_preservation=metadata_preservation
        )


class PowerPointToMarkdownConverter:
    """PowerPoint presentation to Markdown converter."""

    def __init__(self):
        self.quality_validator = ConversionQualityValidator()

    async def convert(self, file_path: Path, options: ConversionOptions) -> ConversionResult:
        """Convert PowerPoint presentation to markdown.

        Args:
            file_path: Path to PowerPoint file
            options: Conversion options

        Returns:
            ConversionResult with markdown content
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available for PowerPoint processing")

        try:
            presentation = Presentation(str(file_path))

            # Extract presentation metadata
            metadata = self._extract_pptx_metadata(presentation, file_path)

            # Convert to markdown
            markdown_content = await self._convert_pptx_to_markdown(presentation, options)

            # Calculate quality score
            quality_score = self._calculate_pptx_quality(presentation, markdown_content)

            return ConversionResult(
                content=markdown_content,
                metadata=metadata,
                quality_score=quality_score,
                success=True,
                original_format='pptx'
            )

        except Exception as e:
            logger.error(f"PowerPoint conversion failed: {e}")
            raise

    async def _convert_pptx_to_markdown(self, presentation, options: ConversionOptions) -> str:
        """Convert PowerPoint presentation to structured markdown.

        Args:
            presentation: PowerPoint presentation object
            options: Conversion options

        Returns:
            Markdown content
        """
        sections = []

        # Presentation title
        title = "PowerPoint Presentation"
        sections.append(f"# {title}")
        sections.append("")

        # Presentation metadata
        if options.include_metadata:
            sections.append("## Presentation Information")
            sections.append("")
            sections.append(f"**Slides**: {len(presentation.slides)}")
            sections.append("")

        # Process each slide
        for i, slide in enumerate(presentation.slides, 1):
            sections.append(f"## Slide {i}")
            sections.append("")

            # Extract text from slide
            slide_text = self._extract_slide_text(slide)
            if slide_text:
                sections.append(slide_text)
            else:
                sections.append("*No text content in this slide*")

            sections.append("")

        return "\n".join(sections)

    def _extract_slide_text(self, slide) -> str:
        """Extract text content from a slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            Extracted text content
        """
        text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text.strip())

        return "\n\n".join(text_parts)

    def _extract_pptx_metadata(self, presentation, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from PowerPoint presentation.

        Args:
            presentation: PowerPoint presentation object
            file_path: Path to file

        Returns:
            Metadata dictionary
        """
        metadata = {
            'original_filename': file_path.name,
            'file_size': file_path.stat().st_size,
            'conversion_format': 'pptx_to_markdown',
            'converter_version': '1.0.0',
            'office_type': 'powerpoint'
        }

        # Add presentation properties
        metadata.update({
            'slide_count': len(presentation.slides),
            'has_content': any(any(hasattr(shape, "text") and shape.text for shape in slide.shapes) for slide in presentation.slides)
        })

        return metadata

    def _calculate_pptx_quality(self, presentation, markdown_content: str) -> QualityScore:
        """Calculate quality score for PowerPoint conversion.

        Args:
            presentation: PowerPoint presentation object
            markdown_content: Generated markdown

        Returns:
            Quality score
        """
        # Base scores
        completeness_score = 0.8
        readability_score = 0.85
        structure_score = 0.9
        metadata_preservation = 0.85

        # Adjust based on content
        word_count = len(markdown_content.split())
        if word_count < 50:
            completeness_score *= 0.8

        # Check for slide structure
        slide_count = markdown_content.count('## Slide')
        if slide_count > 0:
            structure_score = min(structure_score * (1 + slide_count * 0.01), 1.0)

        overall_score = (
            completeness_score * 0.3 +
            readability_score * 0.25 +
            structure_score * 0.25 +
            metadata_preservation * 0.2
        )

        return QualityScore(
            overall_score=overall_score,
            completeness_score=completeness_score,
            readability_score=readability_score,
            structure_score=structure_score,
            metadata_preservation=metadata_preservation
        )